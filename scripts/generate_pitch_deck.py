#!/usr/bin/env python3
"""
Generate a pitch deck (.pptx) styled like frontend/src/index.css (Riscon dashboard).
Run: pip install python-pptx && python3 scripts/generate_pitch_deck.py
Output: pitch-deck/CursorHackathon-Pitch.pptx

Slide flow: cover + 5 investor slides (market, solution, business model, GTM, team) + Q&A.
Replace bracketed [EDIT: …] placeholders with your sourced numbers and names.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

# --- Theme from frontend :root (index.css) ---
BG_PAGE = RGBColor(0x0A, 0x0F, 0x1C)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)
TEXT_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_TERTIARY = RGBColor(0x64, 0x74, 0x8B)
SUCCESS = RGBColor(0x22, 0xC5, 0x5E)
WARNING = RGBColor(0xFA, 0xCC, 0x15)
BORDER = RGBColor(0x1E, 0x29, 0x3B)

W = Inches(13.333)  # 16:9
H = Inches(7.5)
M = Inches(0.55)
TOP_BAR = Inches(0.12)


def add_full_bleed_bg(slide, color: RGBColor) -> None:
    """First shape added sits behind later shapes."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_accent_bar(slide) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, TOP_BAR)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_card(slide, left, top, width, height) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)


def textbox(slide, left, top, width, height, text: str, *, size=14, bold=False, color=TEXT_PRIMARY, font="Calibri", align=PP_ALIGN.LEFT, small_caps=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    if small_caps:
        try:
            r.font.small_caps = True
        except AttributeError:
            pass
    return box


def bullet_slide(
    prs,
    title: str,
    lines: list[str],
    kicker: str | None = None,
    *,
    body_pt: int = 14,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_bg(slide, BG_PAGE)
    add_accent_bar(slide)

    if kicker:
        textbox(slide, M, Inches(0.45), W - 2 * M, Inches(0.35), kicker.upper(), size=10, color=ACCENT, bold=True, font="Consolas", small_caps=True)

    textbox(slide, M, Inches(0.85), W - 2 * M, Inches(0.9), title.upper(), size=26, bold=True, color=TEXT_PRIMARY, font="Calibri Light")

    body_top = Inches(1.85)
    body_h = H - body_top - M
    add_card(slide, M, body_top, W - 2 * M, body_h)

    tb = slide.shapes.add_textbox(M + Inches(0.35), body_top + Inches(0.35), W - 2 * M - Inches(0.7), body_h - Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(body_pt)
        p.font.color.rgb = TEXT_SECONDARY
        p.font.name = "Calibri"
    return slide


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_bg(slide, BG_PAGE)
    add_accent_bar(slide)

    textbox(
        slide,
        M,
        Inches(2.1),
        W - 2 * M,
        Inches(1.1),
        "RISCON",
        size=54,
        bold=True,
        color=TEXT_PRIMARY,
        font="Calibri Light",
        align=PP_ALIGN.LEFT,
    )
    textbox(
        slide,
        M,
        Inches(3.15),
        W - 2 * M,
        Inches(0.55),
        "SUPPLY CHAIN RISK COMMAND CENTER",
        size=14,
        bold=True,
        color=ACCENT,
        font="Consolas",
        align=PP_ALIGN.LEFT,
        small_caps=True,
    )
    textbox(
        slide,
        M,
        Inches(3.85),
        W - 2 * M,
        Inches(1.2),
        "AI-classified maritime & logistics news, fused with enterprise plants, suppliers, and fleet mobility — delivered through a live operations dashboard and probability signals.",
        size=16,
        color=TEXT_SECONDARY,
        font="Calibri",
        align=PP_ALIGN.LEFT,
    )

    # KPI-style chips (like dashboard)
    chip_y = Inches(5.15)
    chip_w = Inches(2.85)
    chip_h = Inches(0.65)
    gap = Inches(0.35)
    labels = [
        ("LIVE SIGNALS", SUCCESS),
        ("PROBABILITY ENGINE", ACCENT),
        ("ENTERPRISE DATA", WARNING),
    ]
    x = M
    for label, c in labels:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, chip_y, chip_w, chip_h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = BG_CARD
        sh.line.color.rgb = BORDER
        t = slide.shapes.add_textbox(x + Inches(0.2), chip_y + Inches(0.14), chip_w - Inches(0.35), chip_h)
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.name = "Consolas"
        p.font.color.rgb = c
        x += chip_w + gap

    textbox(slide, M, H - Inches(0.65), W - 2 * M, Inches(0.4), "CursorHackathon · Vector-Devs", size=11, color=TEXT_TERTIARY, font="Calibri")
    return slide


def main():
    out_dir = Path(__file__).resolve().parent.parent / "pitch-deck"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "CursorHackathon-Pitch.pptx"

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Cover
    title_slide(prs)

    # Slide 1 — The Market Problem (investors fund markets, not features)
    bullet_slide(
        prs,
        "The market problem",
        [
            "WHO HAS THIS PROBLEM? Heads of supply chain planning and logistics control-tower leads at mid-market and large manufacturers ($100M–$5B revenue), plus ocean-freight and 3PL operations teams serving CPG, automotive, and industrial OEMs. Secondary buyers: procurement risk managers and resilience program owners replacing spreadsheet war rooms.",
            "HOW BIG IS THE MARKET? [EDIT: add SAM — e.g. ~X,000 target accounts in EU+US with offshore suppliers] · [EDIT: attach $B TAM for supply-chain risk / control-tower software with source] · [EDIT: user seats or ARR pool you believe you can capture in 5 years]",
            "HOW DO THEY SOLVE IT TODAY, AND WHY DOES THAT SUCK? Email chains, Google Alerts, shared Excel risk registers, and generic news dashboards. None of it is wired to their plants, suppliers, active shipments, or fleet positions, so leaders get noise, 24–48h lag, and no shared probability of impact. Enterprise GRC suites are heavy, slow to deploy, and tuned for audit—not daily ops.",
            "Investors fund markets, not features. We are going after the gap between “headline chaos” and “actionable, portfolio-level risk.”",
        ],
        kicker="Slide 1",
        body_pt=13,
    )

    # Slide 2 — Your Solution
    bullet_slide(
        prs,
        "Your solution",
        [
            "ONE SENTENCE: Riscon is an AI-native supply-chain risk cockpit that turns maritime and logistics news into structured signals, scores them against your enterprise footprint and fleet context, and pushes live probability updates to operators.",
            "WHY ~10× BETTER THAN ALTERNATIVES? Versus alerts + spreadsheets: same-day structured signals tied to sites and suppliers, not detached headlines. Versus legacy resilience suites: faster time-to-value (Docker-ready microservices, API-first), focused UX for the control tower, and continuous scoring instead of quarterly manual assessments. Versus generic LLM chat: productized pipelines, caching, and WebSocket delivery built for ops, not ad-hoc prompts.",
            "Keep the demo tight: one screen that shows news → probability → map. That is enough to “get it.”",
        ],
        kicker="Slide 2",
    )

    # Slide 3 — Business Model
    bullet_slide(
        prs,
        "Business model",
        [
            "WHO PAYS? VP / Director Supply Chain, Chief Procurement Officer, Head of Logistics, and increasingly Chief Risk (operational resilience budget). Second line: 3PLs and 4PLs who resell “white-label” risk visibility to shippers.",
            "HOW MUCH WOULD THEY PAY? [EDIT: pick anchor] Example framing: $25–75 per named user / month for the team tier, or $18k–60k annual platform fee for mid-market (10–50 sites) with minimums. Enterprise above that with SSO, SLAs, and custom data feeds.",
            "UNIT ECONOMICS (ILLUSTRATIVE — REPLACE WITH YOUR MODEL): Assume $40k ACV, 75% gross margin after cloud + LLM API usage. If news enrichment costs ~$[EDIT] per active customer per month and hosting ~$[EDIT], target CAC payback in <12 months via inside sales. Show consumption caps on AI to protect margin.",
        ],
        kicker="Slide 3",
        body_pt=13,
    )

    # Slide 4 — Go-to-Market
    bullet_slide(
        prs,
        "Go-to-market plan",
        [
            "FIRST 10 CUSTOMERS (BE SPECIFIC — REPLACE WITH YOUR LIST): (1) DM or warm intro to 30 heads of logistics in [EDIT: e.g. German automotive supplier association members]. (2) Offer free 30-day pilot to 5 mid-market manufacturers you can reach via [EDIT: e.g. IHK / VDMA events]. (3) Partner pitch to 2 regional 3PLs for a co-branded “shipper risk” widget. (4) Publish one public case: “From alert to action in 15 minutes.”",
            "DISTRIBUTION ADVANTAGE: [EDIT — e.g. existing enterprise relationships, open-source community, alumni network, niche newsletter, or integration into WMS/TMS partner marketplace]. What is repeatable beyond “posting on LinkedIn”?",
            "WHY SWITCH FROM THE STATUS QUO? They leave when weekly executive reviews still start with “did anyone see the news about…” Switch when Riscon shows probability by site/supplier on a live board their CEO can trust. Offer side-by-side pilot vs. current alert inbox for 4 weeks.",
        ],
        kicker="Slide 4",
    )

    # Slide 5 — Why You / Why Now
    bullet_slide(
        prs,
        "Why you · why now",
        [
            "WHY YOU? [EDIT: team bullets — domain years in supply chain / logistics, shipped ML or data products, enterprise sales experience, design for ops]. Vector-Devs built an end-to-stack reference: Spring Boot agents, probability service, React operations UI, Docker compose — proof you can ship, not only slide.",
            "FULL-TIME IF IT TRACTIONS? [EDIT: Yes / timeline — e.g. “All founders commit to full-time on $XM ARR or institutional seed.”]",
            "WHY NOW? Geopolitical volatility, near-shoring pressure, and insurer / board demand for operational resilience are expanding budgets. LLM costs and APIs finally allow affordable per-article classification at scale. Incumbents are not optimized for real-time, map-first workflows.",
        ],
        kicker="Slide 5",
        body_pt=13,
    )

    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_bg(slide, BG_PAGE)
    add_accent_bar(slide)
    textbox(slide, M, Inches(2.9), W - 2 * M, Inches(1.2), "THANK YOU", size=44, bold=True, color=TEXT_PRIMARY, font="Calibri Light", align=PP_ALIGN.CENTER)
    textbox(slide, M, Inches(4.15), W - 2 * M, Inches(0.6), "Q & A", size=20, color=ACCENT, font="Consolas", align=PP_ALIGN.CENTER, bold=True)
    textbox(slide, M, Inches(5.0), W - 2 * M, Inches(0.8), "github.com/Vector-Devs/CursorHackathon", size=14, color=TEXT_SECONDARY, font="Calibri", align=PP_ALIGN.CENTER)

    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
